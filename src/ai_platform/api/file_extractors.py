"""File content extraction for chat context injection."""

import csv
import io
from pathlib import Path

import structlog

logger = structlog.get_logger(__name__)

MAX_CHARS = 8000

TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".json",
    ".xml",
    ".yaml",
    ".yml",
    ".py",
    ".js",
    ".ts",
    ".html",
    ".css",
    ".sql",
    ".sh",
    ".log",
}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp"}


def extract_file_content(fpath: Path) -> str | None:
    """Extract text content from a file. Returns text or None if extraction fails."""
    ext = fpath.suffix.lower()

    if ext in IMAGE_EXTENSIONS:
        return None

    if ext == ".csv":
        return _extract_csv(fpath)

    if ext in TEXT_EXTENSIONS:
        return _extract_text(fpath)

    extractor = _EXTRACTORS.get(ext)
    if extractor:
        return extractor(fpath)

    return None


def _extract_text(fpath: Path) -> str | None:
    try:
        raw = fpath.read_bytes()[:MAX_CHARS]
        try:
            text = raw.decode("utf-8")
        except UnicodeDecodeError:
            try:
                text = raw.decode("latin-1")
            except Exception:
                return None
        if text.count("\ufffd") > len(text) * 0.1:
            return None
        return text
    except Exception:
        return None


def _extract_csv(fpath: Path) -> str | None:
    try:
        raw = fpath.read_bytes()[: MAX_CHARS * 2]
        text = raw.decode("utf-8", errors="replace")
        reader = csv.reader(io.StringIO(text))
        rows = list(reader)
        if not rows:
            return "(empty CSV file)"
        header = rows[0]
        lines = ["| " + " | ".join(header) + " |"]
        lines.append("| " + " | ".join(["---"] * len(header)) + " |")
        for row in rows[1:]:
            line = "| " + " | ".join(row) + " |"
            lines.append(line)
            if sum(len(ln) for ln in lines) > MAX_CHARS:
                lines.append(f"... (truncated, {len(rows)} total rows)")
                break
        return "\n".join(lines)[:MAX_CHARS]
    except Exception:
        return None


def _extract_pdf(fpath: Path) -> str | None:
    try:
        import fitz

        doc = fitz.open(str(fpath))
        text = ""
        for page in doc:
            text += page.get_text()
            if len(text) > MAX_CHARS:
                break
        doc.close()
        return text[:MAX_CHARS] if text.strip() else "(empty PDF)"
    except Exception:
        return None


def _extract_docx(fpath: Path) -> str | None:
    try:
        from docx import Document

        doc = Document(str(fpath))
        text = "\n".join(p.text for p in doc.paragraphs)
        return text[:MAX_CHARS] if text.strip() else "(empty document)"
    except Exception:
        return None


def _extract_xlsx(fpath: Path) -> str | None:
    try:
        from openpyxl import load_workbook

        wb = load_workbook(str(fpath), read_only=True, data_only=True)
        parts: list[str] = []
        for sheet in wb.worksheets:
            parts.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                line = " | ".join(str(c) if c is not None else "" for c in row)
                parts.append(line)
                if sum(len(p) for p in parts) > MAX_CHARS:
                    parts.append("... (truncated)")
                    break
            if sum(len(p) for p in parts) > MAX_CHARS:
                break
        wb.close()
        return "\n".join(parts)[:MAX_CHARS]
    except Exception:
        return None


def _extract_pptx(fpath: Path) -> str | None:
    try:
        from pptx import Presentation

        prs = Presentation(str(fpath))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_texts.append(shape.text_frame.text)
            if slide_texts:
                parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
            if sum(len(p) for p in parts) > MAX_CHARS:
                parts.append("... (truncated)")
                break
        return "\n\n".join(parts)[:MAX_CHARS] if parts else "(empty presentation)"
    except Exception:
        return None


def _extract_odt(fpath: Path) -> str | None:
    try:
        from odf import teletype
        from odf.opendocument import load
        from odf.text import P

        doc = load(str(fpath))
        paragraphs = doc.getElementsByType(P)
        text = "\n".join(teletype.extractText(p) for p in paragraphs)
        return text[:MAX_CHARS] if text.strip() else "(empty document)"
    except Exception:
        return None


def _extract_ods(fpath: Path) -> str | None:
    try:
        from odf import teletype
        from odf.opendocument import load
        from odf.table import Table, TableCell, TableRow

        doc = load(str(fpath))
        parts: list[str] = []
        for table in doc.getElementsByType(Table):
            parts.append(f"[Sheet: {table.getAttribute('name')}]")
            for row in table.getElementsByType(TableRow):
                cells = [teletype.extractText(c) for c in row.getElementsByType(TableCell)]
                parts.append(" | ".join(cells))
                if sum(len(p) for p in parts) > MAX_CHARS:
                    parts.append("... (truncated)")
                    break
            if sum(len(p) for p in parts) > MAX_CHARS:
                break
        return "\n".join(parts)[:MAX_CHARS]
    except Exception:
        return None


def _extract_rtf(fpath: Path) -> str | None:
    try:
        from striprtf.striprtf import rtf_to_text

        raw = fpath.read_text(errors="replace")
        text = rtf_to_text(raw)
        return text[:MAX_CHARS] if text.strip() else "(empty document)"
    except Exception:
        return None


def _extract_epub(fpath: Path) -> str | None:
    try:
        import ebooklib
        from bs4 import BeautifulSoup
        from ebooklib import epub

        book = epub.read_epub(str(fpath), options={"ignore_ncx": True})
        parts: list[str] = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_body_content(), "html.parser")
            text = soup.get_text(separator="\n", strip=True)
            if text:
                parts.append(text)
            if sum(len(p) for p in parts) > MAX_CHARS:
                break
        return "\n\n".join(parts)[:MAX_CHARS] if parts else "(empty epub)"
    except Exception:
        return None


def _extract_doc(fpath: Path) -> str | None:
    try:
        import fitz

        doc = fitz.open(str(fpath))
        text = ""
        for page in doc:
            text += page.get_text()
            if len(text) > MAX_CHARS:
                break
        doc.close()
        return text[:MAX_CHARS] if text.strip() else None
    except Exception:
        return None


_EXTRACTORS: dict[str, callable] = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".doc": _extract_doc,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
    ".odt": _extract_odt,
    ".ods": _extract_ods,
    ".rtf": _extract_rtf,
    ".epub": _extract_epub,
}
